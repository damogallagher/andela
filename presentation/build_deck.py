#!/usr/bin/env python3
"""Build the Andela Guardrail Auditor presentation (.pptx).

Redesigned deck: act-divider structure with a forest-green brand system,
ghost numerals, numbered badges, framed product screenshots and stat callouts.
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ASSETS = Path(__file__).resolve().parent / "assets"
OUT = Path(__file__).resolve().parent / "Andela-Guardrail-Auditor.pptx"

# ---- palette (forest / moss brand) ----------------------------------------
FOREST = RGBColor(0x0F, 0x3D, 0x2E)   # dark green
DEEP = RGBColor(0x07, 0x24, 0x1B)     # near-black green (gradient base)
GHOST = RGBColor(0x1B, 0x55, 0x41)    # subtle ghost numeral on dark
GREEN = RGBColor(0x1F, 0x9D, 0x6B)    # accent
MOSS = RGBColor(0x2C, 0x5F, 0x2D)
INK = RGBColor(0x1B, 0x2A, 0x24)      # body on light
MUTE = RGBColor(0x5A, 0x6B, 0x63)     # captions
FAINT = RGBColor(0x9A, 0xA8, 0xA2)    # page numbers
PANEL = RGBColor(0xEE, 0xF4, 0xF1)    # light tint card
PANEL2 = RGBColor(0xDD, 0xE8, 0xE2)   # image border
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xD9, 0xF2, 0xE7)      # light text on dark
CRIT = RGBColor(0xC0, 0x14, 0x3C)
HIGH = RGBColor(0xD9, 0x62, 0x2B)
MED = RGBColor(0xC8, 0x88, 0x1A)
LOW = RGBColor(0x1F, 0x9D, 0x6B)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"
MONO_FONT = "Courier New"

SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

PAGE = {"n": 0}


def _rm(tf):
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0


def text(s, l, t, w, h, runs, size=16, color=INK, bold=False, font=BODY_FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=4):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    _rm(tf)
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = list(runs.split("\n"))
    norm = []
    for line in runs:
        if isinstance(line, str):
            norm.append([(line, {})])
        else:
            segs = []
            for seg in line:
                segs.append((seg, {}) if isinstance(seg, str) else (seg[0], seg[1]))
            norm.append(segs)
    first = True
    for line in norm:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for seg, opt in line:
            r = p.add_run()
            r.text = seg
            r.font.name = opt.get("font", font)
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.italic = opt.get("italic", False)
            r.font.color.rgb = opt.get("color", color)
    return tb


def _shadow(sp, alpha="28000", blur="90000", dist="38000"):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn("a:effectLst"), {})
    sh = spPr.makeelement(qn("a:outerShdw"),
                          {"blurRad": blur, "dist": dist, "dir": "5400000", "rotWithShape": "0"})
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": "0B2A20"})
    a = spPr.makeelement(qn("a:alpha"), {"val": alpha})
    clr.append(a)
    sh.append(clr)
    el.append(sh)
    spPr.append(el)


def _grad(sp, c1, c2, angle=90):
    """Apply a linear gradient (angle deg: 90 = top->bottom)."""
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = spPr.makeelement(qn("a:gsLst"), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = spPr.makeelement(qn("a:gs"), {"pos": str(pos)})
        clr = spPr.makeelement(qn("a:srgbClr"), {"val": str(col)})
        gs.append(clr)
        gsLst.append(gs)
    grad.append(gsLst)
    lin = spPr.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000)), "scaled": "1"})
    grad.append(lin)
    ln = spPr.find(qn("a:ln"))
    (ln.addprevious(grad) if ln is not None else spPr.append(grad))


def slide(bg=WHITE, grad_to=None):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    if grad_to is not None:
        _grad(r, bg, grad_to, angle=115)
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.08, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        _shadow(sp)
    return sp


def picture(s, path, l, t, w=None, h=None, shadow=True, border=PANEL2):
    iw, ih = Image.open(path).size
    ratio = iw / ih
    if w and not h:
        h = w / ratio
    elif h and not w:
        w = h * ratio
    pic = s.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
    pic.line.color.rgb = border
    pic.line.width = Pt(1)
    pic.shadow.inherit = False
    if shadow:
        _shadow(pic)
    return pic, w, h


def circle(s, l, t, d, fill, glyph="", gsize=16, gcolor=WHITE, font=BODY_FONT):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    c.shadow.inherit = False
    if glyph:
        tf = c.text_frame
        tf.word_wrap = False
        _rm(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = glyph
        r.font.name = font
        r.font.size = Pt(gsize)
        r.font.bold = True
        r.font.color.rgb = gcolor
    return c


def pagenum(s, dark=False):
    PAGE["n"] += 1
    col = ICE if dark else FAINT
    text(s, 11.1, 7.04, 1.9, 0.3,
         [[("Guardrail Auditor   ", {"color": col, "size": 9.5}),
           ("· %02d" % PAGE["n"], {"color": GREEN if not dark else GREEN, "size": 9.5, "bold": True})]],
         align=PP_ALIGN.RIGHT)


def header(s, kicker, title, section=""):
    text(s, 0.7, 0.5, 9.5, 0.3, kicker.upper(), size=13.5, color=GREEN, bold=True,
         font=BODY_FONT)
    text(s, 0.7, 0.8, 11.8, 0.85, title, size=33, color=FOREST, bold=True, font=HEAD_FONT)
    if section:
        text(s, 9.6, 0.52, 3.05, 0.3, section.upper(), size=10.5, color=FAINT, bold=True,
             align=PP_ALIGN.RIGHT)
    pagenum(s)


def divider(num, act, title, subtitle):
    s = slide(FOREST, grad_to=DEEP)
    # ghost numeral
    text(s, 7.2, 0.3, 6.0, 7.2, num, size=420, color=GHOST, bold=True, font=HEAD_FONT,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # accent dot motif
    circle(s, 0.95, 3.02, 0.16, GREEN)
    text(s, 1.3, 2.9, 8.0, 0.4, act.upper(), size=15, color=GREEN, bold=True)
    text(s, 0.9, 3.35, 8.2, 1.5, title, size=46, color=WHITE, bold=True, font=HEAD_FONT,
         line_spacing=1.0)
    text(s, 0.95, 5.1, 7.6, 1.0, subtitle, size=18, color=ICE, line_spacing=1.2)
    return s


# ===========================================================================
# COVER
# ===========================================================================
s = slide(FOREST, grad_to=DEEP)
for d, x, y in [(4.2, 10.0, -1.4), (2.6, 12.0, 4.3)]:
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.background()
    o.line.color.rgb = GHOST
    o.line.width = Pt(2)
    o.shadow.inherit = False
circle(s, 0.95, 1.42, 0.16, GREEN)
text(s, 1.3, 1.3, 11.0, 0.4, "ANDELA GRADUATE “VIBE CODING” CHALLENGE  ·  PROJECT 2",
     size=14.5, color=GREEN, bold=True)
text(s, 0.9, 2.0, 11.6, 2.0, "Enterprise Security\nGuardrail Auditor",
     size=58, color=WHITE, bold=True, font=HEAD_FONT, line_spacing=0.98)
text(s, 0.95, 4.35, 11.0, 0.6,
     "Shift-left security scanning for Infrastructure-as-Code", size=23, color=ICE,
     font=HEAD_FONT)
pill = rect(s, 0.95, 5.35, 5.75, 0.6, fill=GREEN, radius=0.5)
text(s, 0.95, 5.35, 5.75, 0.6, "Tagle Tag  ·  Connector — Foundation Operator",
     size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.95, 6.55, 11.6, 0.4,
     [[("Damien Gallagher", {"bold": True, "color": WHITE}),
       ("       github.com/damogallagher/andela", {"color": ICE})]], size=15)

# ===========================================================================
# ACT 01 — THE MANDATE
# ===========================================================================
divider("01", "Act one", "The mandate",
        "A “Vibe Coding” exercise: I hold the architectural vision, the AI agent writes every line.")

# Slide — I architect, the AI engineers
s = slide()
header(s, "The challenge", "I architect, the AI engineers", "01 · The mandate")
rows = [
    ("1", "Human-in-the-loop architect",
     "I set vision, constraints and acceptance criteria; the AI agent wrote 100% of the code — no manual edits."),
    ("2", "Auditable by design",
     "Every instruction is logged in prompts.md — 58 prompts forming a complete trail from scaffold to production-ready."),
    ("3", "One agent, end-to-end",
     "The same AI tool throughout, preserving architectural consistency across every layer of the stack."),
    ("4", "Time-boxed delivery",
     "A 4–6 hour MVP target, shipped as a public repo, a full prompt log, and this generated deck."),
]
y = 2.05
for n, h, b in rows:
    circle(s, 0.7, y, 0.52, GREEN, n, gsize=19)
    text(s, 1.45, y - 0.04, 11.2, 0.4, h, size=19, color=FOREST, bold=True)
    text(s, 1.45, y + 0.38, 11.2, 0.6, b, size=14.5, color=MUTE, line_spacing=1.08)
    y += 1.12
rect(s, 0.7, 6.36, 11.95, 0.56, fill=PANEL, radius=0.12)
text(s, 1.0, 6.36, 11.4, 0.56,
     [[("My value was not syntax — it was ", {"color": INK, "italic": True}),
       ("directing the system: architecture, interfaces, the quality bar, and the trade-offs.",
        {"color": FOREST, "bold": True, "italic": True})]],
     size=15, anchor=MSO_ANCHOR.MIDDLE)

# Slide — the problem
s = slide()
header(s, "Why it matters", "Breaches come from misconfiguration", "01 · The mandate")
text(s, 0.7, 1.95, 5.95, 1.4,
     "Most cloud incidents are not zero-days — they are risky settings that ship in Terraform "
     "and CloudFormation before anything ever reaches the cloud.",
     size=17, color=INK, line_spacing=1.2)
rect(s, 0.7, 5.35, 5.95, 1.35, fill=FOREST, radius=0.08, shadow=True)
text(s, 1.0, 5.55, 5.4, 1.0,
     [[("The cheapest place to catch them is the ", {"color": ICE, "size": 16}),
       ("pull request", {"color": WHITE, "bold": True, "size": 16}),
       (" — not production.", {"color": ICE, "size": 16})]], line_spacing=1.2,
     anchor=MSO_ANCHOR.MIDDLE)
cards = [
    ("Open SSH ingress", "Security groups exposing port 22 to 0.0.0.0/0", CRIT),
    ("Public S3 buckets", "Public ACLs and disabled versioning", HIGH),
    ("Wildcard IAM", "Action:* on Resource:* — full privilege", HIGH),
    ("Hardcoded secrets", "Credentials and AWS keys committed to IaC", CRIT),
    ("Unencrypted databases", "Storage encryption switched off", MED),
    ("No backup / retention", "Deletion protection and versioning suspended", LOW),
]
x0, y0, cw, ch, gx, gy = 6.95, 1.92, 2.86, 1.5, 0.2, 0.22
for i, (t_, b_, col) in enumerate(cards):
    cx = x0 + (i % 2) * (cw + gx)
    cy = y0 + (i // 2) * (ch + gy)
    rect(s, cx, cy, cw, ch, fill=PANEL, radius=0.1, shadow=True)
    circle(s, cx + 0.22, cy + 0.24, 0.22, col)
    text(s, cx + 0.56, cy + 0.2, cw - 0.72, 0.5, t_, size=14, color=FOREST, bold=True)
    text(s, cx + 0.24, cy + 0.72, cw - 0.46, 0.7, b_, size=11.5, color=MUTE, line_spacing=1.05)

# ===========================================================================
# ACT 02 — THE SOLUTION
# ===========================================================================
divider("02", "Act two", "The solution",
        "A Python, API-first guardrail auditor — one scanner core serving humans and pipelines alike.")

# Slide — the solution
s = slide()
header(s, "The solution", "A guardrail auditor for humans and pipelines", "02 · The solution")
text(s, 0.7, 1.85, 12, 0.6,
     [[("A ", {}), ("Python, API-first", {"bold": True, "color": FOREST}),
       (" auditor that scans IaC against a security baseline and presents a visual Risk Score.", {})]],
     size=18, color=INK, line_spacing=1.1)
specs = [
    ("Scans", "Terraform (.tf) + JSON / CloudFormation"),
    ("Detects", "100+ AWS / Azure signatures, 9 control families"),
    ("Scores", "Normalized 0–100 risk health, color-coded"),
    ("Persists", "Postgres-backed scan history & audit trail"),
    ("Surfaces", "Dashboard · REST API · CLI gate · SARIF"),
]
y = 2.7
for k, v in specs:
    rect(s, 0.7, y, 2.45, 0.6, fill=FOREST, radius=0.18)
    text(s, 0.7, y, 2.45, 0.6, k, size=15, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 3.27, y, 5.05, 0.6, fill=PANEL, radius=0.18)
    text(s, 3.52, y, 4.7, 0.6, v, size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.76
rect(s, 8.72, 2.7, 3.93, 4.0, fill=PANEL, radius=0.06, shadow=True)
circle(s, 9.02, 3.0, 0.42, GREEN, "◆", gsize=15)
text(s, 9.58, 3.04, 3.0, 0.5, "Local-first by design", size=17, color=FOREST, bold=True,
     font=HEAD_FONT)
text(s, 9.02, 3.75, 3.35, 2.8,
     "Runs entirely on Docker Compose — no cloud credentials needed to develop or review.\n\n"
     "AWS infrastructure is fully authored in Terraform but never auto-provisions, honoring "
     "the challenge’s no-resources rule.",
     size=14, color=MUTE, line_spacing=1.2)

# Slide — five use cases
s = slide()
header(s, "Who uses it", "Five use cases", "02 · The solution")
uc = [
    ("PR / pre-merge gate",
     "python -m app.cli scan --fail-on critical blocks risky infra in CI before merge."),
    ("Developer self-service",
     "Engineers upload a .tf / .json file and see findings plus fixes instantly in the dashboard."),
    ("GitHub Code Scanning",
     "SARIF export lights up the Security tab and inline pull-request annotations."),
    ("Compliance audit trail",
     "Every scan is persisted in Postgres with timestamps, providing durable evidence."),
    ("Regression detection",
     "Compare two scans to prove a change introduced — or resolved — risk."),
]
y = 2.0
for i, (h, b) in enumerate(uc):
    circle(s, 0.7, y, 0.54, GREEN if i % 2 == 0 else MOSS, str(i + 1), gsize=20)
    text(s, 1.5, y - 0.02, 11.1, 0.4, h, size=18.5, color=FOREST, bold=True)
    text(s, 1.5, y + 0.4, 11.1, 0.5, b, size=14.5, color=MUTE, line_spacing=1.05)
    y += 1.0

# Slide — architecture
s = slide()
header(s, "Architecture", "One scanner core, four interfaces", "02 · The solution")
cons = [("React Dashboard", "upload · filter · trend · SARIF"),
        ("Pipeline CLI", "python -m app.cli scan"),
        ("GitHub Actions", "lint · test · SARIF · gate")]
cy = 2.15
for t_, sub in cons:
    rect(s, 0.7, cy, 3.0, 1.05, fill=PANEL, radius=0.12, shadow=True)
    text(s, 0.9, cy + 0.16, 2.7, 0.4, t_, size=14.5, color=FOREST, bold=True)
    text(s, 0.9, cy + 0.58, 2.7, 0.4, sub, size=11, color=MUTE)
    cy += 1.36
rect(s, 4.6, 2.6, 3.0, 1.55, fill=FOREST, radius=0.1, shadow=True)
text(s, 4.6, 2.78, 3.0, 0.5, "FastAPI app", size=18, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, font=HEAD_FONT)
text(s, 4.75, 3.3, 2.7, 0.7, "REST scans, upload,\ncompare, SARIF, health",
     size=12, color=ICE, align=PP_ALIGN.CENTER, line_spacing=1.05)
rect(s, 4.6, 4.55, 3.0, 1.4, fill=GREEN, radius=0.1, shadow=True)
text(s, 4.6, 4.74, 3.0, 0.5, "Scanner core", size=17, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, font=HEAD_FONT)
text(s, 4.75, 5.24, 2.7, 0.6, "Terraform rules + JSON rules", size=12, color=WHITE,
     align=PP_ALIGN.CENTER, line_spacing=1.05)
rect(s, 8.5, 2.6, 3.05, 1.55, fill=PANEL, radius=0.1, shadow=True)
text(s, 8.5, 2.84, 3.05, 0.5, "Postgres", size=16, color=FOREST, bold=True, align=PP_ALIGN.CENTER)
text(s, 8.65, 3.34, 2.75, 0.6, "scan runs + findings\n(audit trail)", size=12, color=MUTE,
     align=PP_ALIGN.CENTER, line_spacing=1.05)
rect(s, 8.5, 4.55, 3.05, 1.4, fill=PANEL, radius=0.1, shadow=True)
text(s, 8.5, 4.8, 3.05, 0.5, "IaC sources", size=16, color=FOREST, bold=True, align=PP_ALIGN.CENTER)
text(s, 8.65, 5.3, 2.75, 0.5, "sample_iac / uploads", size=12, color=MUTE, align=PP_ALIGN.CENTER)


def connect(x1, y1, x2, y2, col=MOSS):
    cxn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = col
    cxn.line.width = Pt(2.25)
    cxn.shadow.inherit = False


connect(3.7, 2.67, 4.6, 3.05)
connect(3.7, 3.52, 4.6, 3.35)
connect(3.7, 4.83, 4.6, 5.05)
connect(6.1, 4.15, 6.1, 4.55)
connect(7.6, 3.35, 8.5, 3.35)
connect(7.6, 5.2, 8.5, 5.2)
rect(s, 0.7, 6.34, 11.95, 0.6, fill=PANEL, radius=0.12)
text(s, 1.0, 6.34, 11.4, 0.6,
     [[("Rule logic lives once in ", {"color": INK}),
       ("app.scanner.RULES", {"font": MONO_FONT, "color": FOREST, "bold": True}),
       (" — reused by dashboard, API, CLI, tests and SARIF. Five trade-offs recorded as ADRs in docs/adr/.",
        {"color": INK})]], size=13.5, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# ACT 03 — THE PRODUCT
# ===========================================================================
divider("03", "Act three", "The product",
        "A polished React dashboard that makes infrastructure risk legible at a glance.")

# Slide — risk score (red hero)
s = slide()
header(s, "Live application", "The Risk Score, color-coded", "03 · The product")
text(s, 0.7, 1.74, 12, 0.4,
     [[("Thresholds make risk legible instantly:    ", {"color": INK}),
       ("green > 90", {"color": GREEN, "bold": True}),
       ("     ·     ", {"color": FAINT}),
       ("amber 70–90", {"color": MED, "bold": True}),
       ("     ·     ", {"color": FAINT}),
       ("red < 70", {"color": CRIT, "bold": True})]], size=16)
picture(s, ASSETS / "hero-red.png", 3.18, 2.28, h=4.4)
text(s, 0.7, 6.86, 12, 0.4,
     "A high-risk scan: 47%, 67 findings across 10 files — severity breakdown and a regression alert in one view.",
     size=12.5, color=MUTE, align=PP_ALIGN.CENTER)

# Slide — green & amber
s = slide()
header(s, "Live application", "Green & amber score states", "03 · The product")
picture(s, ASSETS / "hero-green.png", 0.7, 2.05, w=6.0)
picture(s, ASSETS / "hero-amber.png", 6.95, 2.05, w=6.0)
text(s, 0.7, 5.95, 6.0, 0.9,
     [[("100% — green.  ", {"color": GREEN, "bold": True, "size": 15}),
       ("A clean baseline with zero findings.", {"color": MUTE, "size": 14})]], line_spacing=1.1)
text(s, 6.95, 5.95, 6.0, 0.9,
     [[("80% — amber.  ", {"color": MED, "bold": True, "size": 15}),
       ("One medium finding; the normalized model never collapses to zero on small scans.",
        {"color": MUTE, "size": 14})]], line_spacing=1.1)

# Slide — findings table
s = slide()
header(s, "Live application", "Findings management", "03 · The product")
picture(s, ASSETS / "03-findings-table.png", 0.7, 2.0, w=7.95)
rect(s, 8.95, 2.0, 3.7, 4.55, fill=PANEL, radius=0.06, shadow=True)
text(s, 9.25, 2.3, 3.15, 0.4, "Every finding carries", size=15, color=FOREST, bold=True,
     font=HEAD_FONT)
items = ["Rule + severity", "Affected resource", "File : line location", "Remediation advice"]
yy = 2.85
for it in items:
    circle(s, 9.25, yy + 0.06, 0.14, GREEN)
    text(s, 9.55, yy, 2.9, 0.4, it, size=13.5, color=INK)
    yy += 0.5
text(s, 9.25, 5.05, 3.15, 1.4,
     "Search, sortable columns, configurable rows and pagination handle large result sets — "
     "67 findings across 14 pages.",
     size=13, color=MUTE, line_spacing=1.15)

# Slide — filter & regression
s = slide()
header(s, "Live application", "Filtering & regression detection", "03 · The product")
picture(s, ASSETS / "05-severity-filter.png", 0.7, 2.05, w=6.5)
picture(s, ASSETS / "04-regression.png", 7.55, 2.05, w=4.95)
text(s, 0.7, 6.4, 6.5, 0.8,
     "Click a severity card to filter, with breadcrumb and a clear action.", size=13.5, color=MUTE)
text(s, 7.55, 6.4, 5.0, 0.8,
     [[("Comparison flags ", {"color": MUTE, "size": 13.5}),
       ("“introduced 26 new criticals.”", {"color": FOREST, "bold": True, "size": 13.5})]],
     line_spacing=1.05)

# Slide — API-first & responsive
s = slide()
header(s, "Live application", "API-first & responsive", "03 · The product")
picture(s, ASSETS / "07-api-docs.png", 0.7, 2.0, h=4.35)
picture(s, ASSETS / "08-mobile.png", 9.85, 2.0, h=4.35)
text(s, 0.7, 6.5, 8.2, 0.5,
     "Auto-generated OpenAPI 3.1 docs for every endpoint: scans, upload, compare, SARIF, rules, health.",
     size=13, color=MUTE)
text(s, 9.85, 6.5, 2.9, 0.5, "Verified on mobile Chromium.", size=12.5, color=MUTE)

# Slide — feature catalogue
s = slide()
header(s, "Capabilities", "Feature catalogue", "03 · The product")
cols = [
    ("Scanning", GREEN, [
        "100+ AWS / Azure signatures across 9 control families",
        "Single rule registry feeds scanner, /api/rules and SARIF — zero drift",
        "Secret redaction before persistence, display, CLI and SARIF",
    ]),
    ("Dashboard", MOSS, [
        "Color-coded score, severity cards and filtering, breadcrumbs",
        "Search, sortable headers, pagination",
        "Clickable history, trend chart, regression compare, uploads",
    ]),
    ("Interfaces", FOREST, [
        "React dashboard for review",
        "REST API (OpenAPI 3.1)",
        "CLI gate for CI pipelines",
        "SARIF 2.1.0 for Code Scanning",
    ]),
]
cw = 3.85
for i, (title, col, lst) in enumerate(cols):
    x = 0.7 + i * (cw + 0.2)
    rect(s, x, 1.95, cw, 4.75, fill=PANEL, radius=0.05, shadow=True)
    rect(s, x, 1.95, cw, 0.78, fill=col, radius=0.05)
    rect(s, x, 2.35, cw, 0.38, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, x, 1.95, cw, 0.78, title, size=18, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
    yy = 3.05
    for it in lst:
        circle(s, x + 0.3, yy + 0.07, 0.13, col)
        text(s, x + 0.58, yy, cw - 0.85, 0.9, it, size=13.5, color=INK, line_spacing=1.05)
        yy += 0.9

# ===========================================================================
# ACT 04 — THE ENGINEERING
# ===========================================================================
divider("04", "Act four", "The engineering",
        "Documented decisions, an automated pipeline, and a tool that is itself fully guarded.")

# Slide — risk scoring model (ADR 0005)
s = slide()
header(s, "Key decision · ADR 0005", "An honest risk-scoring model", "04 · The engineering")
text(s, 0.7, 1.85, 12, 0.9,
     [[("The original ", {"color": INK}),
       ("100 − flat_penalty", {"font": MONO_FONT, "color": CRIT}),
       (" score saturated at 0 (4 vs 10 criticals looked identical) and ignored scan scope.",
        {"color": INK})]], size=16.5, line_spacing=1.15)
rect(s, 0.7, 2.9, 7.0, 3.65, fill=PANEL, radius=0.06, shadow=True)
text(s, 1.0, 3.15, 6.4, 0.4, "Normalized, severity-weighted density", size=16, color=FOREST,
     bold=True, font=HEAD_FONT)
text(s, 1.0, 3.72, 6.4, 2.7,
     [[("Severity weights", {"bold": True, "color": MOSS})],
      [("critical 10 · high 6 · medium 3 · low 1", {"font": MONO_FONT, "size": 13.5, "color": INK})],
      [(" ", {"size": 6})],
      [("Scope units", {"bold": True, "color": MOSS})],
      [("(files × 2) + distinct affected resources", {"font": MONO_FONT, "size": 13.5, "color": INK})],
      [(" ", {"size": 6})],
      [("Score", {"bold": True, "color": MOSS})],
      [("ceil( 100 · 4·units / (4·units + Σweights) )", {"font": MONO_FONT, "size": 13.5, "color": INK})]],
     size=14, line_spacing=1.15, space_after=2)
rect(s, 7.95, 2.9, 4.7, 3.65, fill=FOREST, radius=0.06, shadow=True)
text(s, 8.25, 3.15, 4.15, 3.2,
     [[("Asymptotic, not subtractive", {"bold": True, "color": WHITE, "size": 16})],
      [(" ", {"size": 8})],
      [("More findings keep lowering the score — there is no hard floor.", {"color": ICE, "size": 14})],
      [(" ", {"size": 6})],
      [("Any scan with findings caps at 99; only a clean scan shows 100.", {"color": ICE, "size": 14})],
      [(" ", {"size": 6})],
      [("Aligns exactly with the green / amber / red thresholds.", {"color": ICE, "size": 14})]],
     line_spacing=1.2)

# Slide — ADRs
s = slide()
header(s, "Engineering judgment", "Architecture decisions (ADRs)", "04 · The engineering")
text(s, 0.7, 1.8, 12, 0.4, "Every significant trade-off is documented and version-controlled.",
     size=15, color=MUTE)
adrs = [
    ("0001", "Local-only dev & verification", "No cloud cost or credential risk; reviewers run it free"),
    ("0002", "Postgres over SQLite", "Mirrors the RDS deployment shape; durable audit trail"),
    ("0003", "Focused text / JSON scan, no full HCL parser", "Small, testable, MVP-appropriate — trade-off documented"),
    ("0004", "Expose via API + CLI + Dashboard + SARIF", "One core, four audiences; CI fails without a web server"),
    ("0005", "Normalized weighted risk score", "Honest density that respects scope and severity"),
]
y = 2.4
for num, dec, why in adrs:
    rect(s, 0.7, y, 1.15, 0.74, fill=FOREST, radius=0.16)
    text(s, 0.7, y, 1.15, 0.74, num, size=15, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
    rect(s, 1.95, y, 10.7, 0.74, fill=PANEL, radius=0.08)
    text(s, 2.25, y + 0.11, 10.2, 0.32, dec, size=15, color=FOREST, bold=True)
    text(s, 2.25, y + 0.43, 10.2, 0.3, why, size=12.5, color=MUTE)
    y += 0.86

# Slide — CI/CD
s = slide()
header(s, "Automation", "A seven-stage CI/CD pipeline", "04 · The engineering")
text(s, 0.7, 1.78, 12, 0.4,
     [[(".github/workflows/ci-cd.yml", {"font": MONO_FONT, "color": FOREST, "bold": True, "size": 13.5}),
       ("  runs on every pull request and push to dev.", {"color": MUTE, "size": 13.5})]])
stages = [
    ("1", "Backend", "Ruff lint + Python unit & functional tests on a Postgres service"),
    ("2", "Code scanning", "Generate SARIF → upload to GitHub Code Scanning"),
    ("3", "Guardrail CLI", "Risky fixtures must fail, clean fixtures pass on --fail-on critical"),
    ("4", "Frontend", "ESLint + Vite build + Playwright (desktop & mobile)"),
    ("5", "Terraform", "fmt -check + init + validate"),
    ("6", "Docker build", "Image built only after every gate passes"),
    ("7", "Deploy (gated)", "OIDC → Terraform S3 backend → ECS / RDS / ALB"),
]
positions = [(0.7, 2.35), (0.7, 3.36), (0.7, 4.37), (0.7, 5.38),
             (6.75, 2.35), (6.75, 3.36), (6.75, 4.37)]
cw2 = 5.9
for (num, h, b), (x, yy) in zip(stages, positions):
    rect(s, x, yy, cw2, 0.92, fill=PANEL, radius=0.1, shadow=True)
    circle(s, x + 0.2, yy + 0.21, 0.5, GREEN if num != "7" else MOSS, num, gsize=18)
    text(s, x + 0.92, yy + 0.16, cw2 - 1.1, 0.35, h, size=15, color=FOREST, bold=True)
    text(s, x + 0.92, yy + 0.5, cw2 - 1.1, 0.35, b, size=11.5, color=MUTE)
rect(s, 6.75, 5.38, cw2, 0.92, fill=FOREST, radius=0.1, shadow=True)
text(s, 7.0, 5.5, cw2 - 0.4, 0.7,
     [[("Plus  ", {"color": ICE, "size": 12.5}),
       ("Dependabot", {"color": WHITE, "bold": True, "size": 12.5}),
       (" (grouped weekly PR) and a  ", {"color": ICE, "size": 12.5}),
       ("pre-commit", {"color": WHITE, "bold": True, "size": 12.5}),
       (" gate mirroring CI.", {"color": ICE, "size": 12.5})]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

# Slide — quality / tests & coverage
s = slide()
header(s, "Quality", "The guardrail tool is itself guarded", "04 · The engineering")
stats = [
    ("71", "Python tests passing", "scanner · API · CLI · SARIF · migrations · observability", GREEN),
    ("100%", "Statement coverage", "841 / 841 lines — gated in CI at --fail-under=100", FOREST),
    ("36", "Playwright runs", "18 specs across desktop & mobile Chromium", MOSS),
    ("58", "Logged prompts", "every instruction captured in prompts.md", GREEN),
]
cw3, ch3 = 2.92, 2.45
for i, (big, lab, sub, col) in enumerate(stats):
    x = 0.7 + i * (cw3 + 0.13)
    rect(s, x, 1.95, cw3, ch3, fill=PANEL, radius=0.07, shadow=True)
    text(s, x + 0.25, 2.2, cw3 - 0.5, 1.0, big, size=52, color=col, bold=True, font=HEAD_FONT)
    text(s, x + 0.25, 3.32, cw3 - 0.5, 0.4, lab, size=15.5, color=FOREST, bold=True)
    text(s, x + 0.25, 3.75, cw3 - 0.5, 0.6, sub, size=12, color=MUTE, line_spacing=1.1)
rect(s, 0.7, 4.7, 11.95, 1.95, fill=FOREST, radius=0.06, shadow=True)
text(s, 1.0, 4.95, 7.6, 1.5,
     [[("$ ./scripts/test-coverage.sh", {"font": MONO_FONT, "color": ICE, "size": 13.5})],
      [("TOTAL      841 stmts      0 missed      100%", {"font": MONO_FONT, "color": WHITE, "size": 14, "bold": True})],
      [("71 passed in 2.55s", {"font": MONO_FONT, "color": GREEN, "size": 14, "bold": True})]],
     line_spacing=1.4)
text(s, 8.85, 4.98, 3.6, 1.5,
     "Coverage spans scanner, API, CLI, SARIF, migrations, observability and config — "
     "a 100% gate keeps it there.",
     size=13, color=ICE, line_spacing=1.18)

# Slide — cloud, deploy-ready never auto-applied
s = slide()
header(s, "Cloud posture", "Deploy-ready, never auto-applied", "04 · The engineering")
text(s, 0.7, 1.9, 5.9, 1.2,
     "Full Terraform under terraform/ provisions a realistic production stack — but the agent "
     "workspace never provisions cloud.",
     size=16.5, color=INK, line_spacing=1.2)
infra = [
    "ECS Fargate behind an Application Load Balancer",
    "RDS Postgres for scan history · ECR for images",
    "Secrets Manager, IAM roles, CloudWatch logging",
    "State in S3 + DynamoDB lock; GitHub OIDC — no long-lived keys",
]
yy = 3.45
for it in infra:
    circle(s, 0.7, yy + 0.04, 0.14, GREEN)
    text(s, 1.0, yy, 5.6, 0.5, it, size=14, color=INK, line_spacing=1.05)
    yy += 0.66
rect(s, 6.95, 1.95, 5.7, 4.6, fill=FOREST, radius=0.06, shadow=True)
circle(s, 7.3, 2.32, 0.42, GREEN, "✓", gsize=17)
text(s, 7.85, 2.36, 4.5, 0.5, "Cloud judgment without cloud cost", size=17, color=WHITE,
     bold=True, font=HEAD_FONT)
text(s, 7.3, 3.15, 5.05, 3.2,
     [[("Deployment runs only when explicit repository variables and credentials are configured.",
        {"color": ICE, "size": 15})],
      [(" ", {"size": 8})],
      [("This honors the challenge’s “delete all cloud resources” rule by default — the stack is "
        "demonstrable on paper and one config away from live, with zero standing spend.",
        {"color": ICE, "size": 15})]],
     line_spacing=1.25)

# ===========================================================================
# ACT 05 — THE HORIZON
# ===========================================================================
divider("05", "Act five", "The horizon",
        "Where the architecture goes next — each enhancement already has a clear seam.")

# Slide — future enhancements
s = slide()
header(s, "Future enhancements", "A roadmap with clear seams", "05 · The horizon")
future = [
    ("Deeper Terraform", "Swap focused text rules for a real HCL parser — modules, variables, dynamic blocks."),
    ("LLM fallback", "Call a model only when deterministic rules find nothing; off by default for free, predictable CI."),
    ("Policy-as-code", "OPA / Rego custom org baselines; rule severities as configuration."),
    ("Alerting", "Slack / webhook notifications on regression or critical findings."),
    ("Multi-tenant + authz", "Teams, projects and role-based access control."),
    ("Scheduled scans", "Track risk posture across repositories over time."),
]
cw4, ch4 = 5.9, 1.32
for i, (h, b) in enumerate(future):
    x = 0.7 + (i % 2) * (cw4 + 0.15)
    yy = 2.05 + (i // 2) * (ch4 + 0.2)
    rect(s, x, yy, cw4, ch4, fill=PANEL, radius=0.09, shadow=True)
    circle(s, x + 0.25, yy + 0.25, 0.42, GREEN if i % 2 == 0 else MOSS, "→", gsize=18)
    text(s, x + 0.85, yy + 0.18, cw4 - 1.05, 0.4, h, size=16, color=FOREST, bold=True)
    text(s, x + 0.85, yy + 0.58, cw4 - 1.05, 0.65, b, size=12.5, color=MUTE, line_spacing=1.08)

# ===========================================================================
# CLOSING
# ===========================================================================
s = slide(FOREST, grad_to=DEEP)
o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(4.3), Inches(3.4), Inches(3.4))
o.fill.background()
o.line.color.rgb = GHOST
o.line.width = Pt(2)
o.shadow.inherit = False
circle(s, 0.95, 0.92, 0.16, GREEN)
text(s, 1.3, 0.8, 11.0, 0.4, "IN SUMMARY", size=14.5, color=GREEN, bold=True)
text(s, 0.9, 1.3, 11.6, 1.4,
     "A production-shaped guardrail\nauditor, delivered by directing\nan AI agent end-to-end.",
     size=34, color=WHITE, bold=True, font=HEAD_FONT, line_spacing=1.02)
checklist = [
    "Tagle Tag — Connector, Foundation Operator",
    "Public GitHub repository — all source code",
    "prompts.md — full 58-prompt audit log",
    "This AI-generated deck — PPTX + Markdown",
    "Cloud cleanup — local-only; nothing provisioned",
]
yy = 3.55
for it in checklist:
    circle(s, 0.95, yy + 0.02, 0.28, GREEN, "✓", gsize=13)
    text(s, 1.4, yy, 10.5, 0.4, it, size=15.5, color=ICE)
    yy += 0.52
rect(s, 0.95, 6.5, 8.3, 0.6, fill=GREEN, radius=0.5)
text(s, 0.95, 6.5, 8.3, 0.6,
     "Python · API-first · Postgres · 100% coverage · CI/CD · SARIF · Terraform",
     size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(str(OUT))
print("Saved", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
